from pptx import Presentation
p = Presentation()
slide_layout = p.slide_layouts[0]
slide = p.slides.add_slide(slide_layout)
zag1 = slide.shapes[0]
zag1.text="Mavzu : Pyhonda pptx bilan ishlash"
zag2=slide.shapes[1]
zag2.text="sana 12.01.2024"
p.save('date/test.pptx')